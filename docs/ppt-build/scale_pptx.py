"""Scale all shapes (positions, sizes, font sizes) in a PPTX by a factor.

Usage: python scale_pptx.py <input.pptx> <output.pptx> <scale_factor>
"""
import sys
from pptx import Presentation
from pptx.util import Pt

SCALE = float(sys.argv[3]) if len(sys.argv) > 3 else 1.3
SRC = sys.argv[1]
DST = sys.argv[2]

p = Presentation(SRC)

# Scale canvas
old_w = p.slide_width
old_h = p.slide_height
new_w = int(round(old_w * SCALE))
new_h = int(round(old_h * SCALE))
p.slide_width = new_w
p.slide_height = new_h
print(f'Canvas: {old_w/914400:.3f} x {old_h/914400:.3f} in -> {new_w/914400:.3f} x {new_h/914400:.3f} in')

# Scale shapes
total_shapes = 0
total_text_runs = 0
total_rotations = 0

for slide in p.slides:
    for shape in slide.shapes:
        # Scale position and size (only if defined)
        try:
            if shape.left is not None:
                shape.left = int(round(shape.left * SCALE))
            if shape.top is not None:
                shape.top = int(round(shape.top * SCALE))
            if shape.width is not None:
                shape.width = int(round(shape.width * SCALE))
            if shape.height is not None:
                shape.height = int(round(shape.height * SCALE))
            total_shapes += 1
        except Exception as e:
            print(f'  WARN scaling shape {shape.shape_id}: {e}')

        # Scale font sizes in text frames
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                # Some python-pptx versions require paragraph defRPr too
                if para.font.size is not None:
                    try:
                        para.font.size = Pt(round(para.font.size.pt * SCALE, 2))
                    except Exception:
                        pass
                for run in para.runs:
                    if run.font.size is not None:
                        try:
                            run.font.size = Pt(round(run.font.size.pt * SCALE, 2))
                            total_text_runs += 1
                        except Exception:
                            pass

        # Rotation is in degrees (angle). Not affected by uniform scale.
        # But if there's any rotation-based dimension (no), leave alone.

p.save(DST)
print(f'Scaled {total_shapes} shapes, {total_text_runs} text runs by {SCALE}x')
print(f'Saved: {DST}')
