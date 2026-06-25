"""Fix L2 card name/desc overlap in scaled PPTX.

The L2 cards have name and desc text shapes with too-tight vertical spacing.
After uniform scaling, name text overflows into desc area. This script moves
the desc text shape down by 0.04" and reduces its height by 0.04" to create
a clean gap.
"""
import sys
sys.stdout = open(sys.stdout.fileno(), mode='w', encoding='utf-8', buffering=1)
from pptx import Presentation
from pptx.util import Emu

SRC = sys.argv[1]
DST = sys.argv[2]
SHIFT_INCHES = 0.05  # how much to push desc down

p = Presentation(SRC)
slide = p.slides[0]
shift_emu = int(SHIFT_INCHES * 914400)

# Identify L2 card text shapes. The desc shapes have small height (~0.08-0.13 in scaled)
# and are positioned within L2 column area. We can identify them by their y range.
# L2 columns are roughly y=3.16 to 4.85 in 1.3x. Desc text shapes within this range.
# We look for shapes that have small height AND are in the L2 region.
L2_Y_MIN = 3.5  # 1.3x scaled
L2_Y_MAX = 5.0  # 1.3x scaled
L2_DESC_HEIGHT_MAX = 0.15  # desc shapes are short

fixed = 0
for shape in slide.shapes:
    if not shape.has_text_frame or not shape.text_frame.text.strip():
        continue
    y = shape.top / 914400 if shape.top else 0
    h = shape.height / 914400 if shape.height else 0
    if not (L2_Y_MIN <= y <= L2_Y_MAX):
        continue
    if h > L2_DESC_HEIGHT_MAX:
        continue
    # This is an L2 desc text shape - push down and shrink
    shape.top = shape.top + shift_emu
    shape.height = max(1, shape.height - shift_emu)
    fixed += 1

p.save(DST)
print(f'Fixed {fixed} L2 desc text shapes. Saved: {DST}')
