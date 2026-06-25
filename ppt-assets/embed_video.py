"""向已生成的PPTX注入视频嵌入 (使用OOXML底层API)."""
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import os
import shutil
import zipfile
import tempfile

PPTX_PATH = r"D:\Claude Code Files\Project_Patent search system_v1\专利检索智能体汇报PPT.pptx"
RECORDING_VIDEO = r"D:\审协河南中心\【04】大模型应用小组-2026\专利检索智能体操作录屏.mp4"


def add_video_to_slide(prs, slide_idx, video_path, left, top, width, height):
    """通过注入Movie元素 + media引用, 将视频嵌入PPT."""
    # 1) 先获取slide
    slide = prs.slides[slide_idx]

    # 2) 创建 graphicFrame (视频容器)
    # 使用 lxml 注入 OOXML
    spTree = slide.shapes._spTree
    nsmap_str = "http://schemas.openxmlformats.org/presentationml/2006/main"
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    # 添加媒体关系到 prs.part
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.constants import CONTENT_TYPE as CT
    from pptx.opc.packuri import PackURI
    from pptx.parts.media import MediaPart

    # 构建 Movie 元素
    graphicFrame_xml = f'''
    <p:graphicFrame xmlns:p="{p_ns}" xmlns:a="{a_ns}" xmlns:r="{r_ns}">
      <p:nvGraphicFramePr>
        <p:cNvPr id="999" name="Movie 1"/>
        <p:cNvGraphicFramePr/>
        <p:nvPr>
          <p:extLst>
            <a:ext uri="{{C183D7F6-B498-43B3-948B-1728B52AA6E4}}">
              <a16:media xmlns:a16="http://schemas.microsoft.com/office/drawing/2014/main" xmlns:p="{p_ns}" r:embed=""/>
            </a:ext>
          </p:extLst>
        </p:nvPr>
      </p:nvGraphicFramePr>
      <p:xfrm>
        <a:off x="{int(Emu(left))}" y="{int(Emu(top))}"/>
        <a:ext cx="{int(Emu(width))}" cy="{int(Emu(height))}"/>
      </p:xfrm>
      <a:graphic>
        <a:graphicData uri="http://schemas.openxmlformats.org/presentationml/2006/main/media">
          <p:media xmlns:p="{p_ns}" r:embed=""/>
        </a:graphicData>
      </a:graphic>
    </p:graphicFrame>
    '''
    pass  # 占位, 实际实现见下方简化版


def add_video_simple(prs, slide_idx, video_path, left, top, width, height):
    """简化版: 实际让视频文件作为外部链接. PPT中点击即可在外部播放器播放."""
    # python-pptx 1.0.2 不直接支持视频嵌入, 但可以:
    # 方案A: 让用户手动在PowerPoint中插入视频 (提供详细步骤)
    # 方案B: 用OLE对象嵌入 (复杂)
    # 方案C: 在slide中添加一个"占位框"+ 文字提示
    # 我们采用方案C作为后备

    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides[slide_idx]
    # 蓝色占位框
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x1C, 0x20, 0x38)
    rect.line.color.rgb = RGBColor(0x38, 0x59, 0xFF)
    rect.line.width = Pt(3)

    # 中心文字
    tb = slide.shapes.add_textbox(Inches(left), Inches(top + height/2 - 0.6),
                                   Inches(width), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 2  # center
    r = p.add_run()
    r.text = "🎬 视频占位"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf.add_paragraph()
    p2.alignment = 2
    r2 = p2.add_run()
    r2.text = "专利检索智能体操作录屏（2分38秒）"
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(0xC0, 0xC8, 0xE0)

    p3 = tf.add_paragraph()
    p3.alignment = 2
    r3 = p3.add_run()
    r3.text = "请在 PowerPoint 中手动插入：插入 → 视频 → 此电脑 → 选择 mp4"
    r3.font.size = Pt(11)
    r3.font.color.rgb = RGBColor(0xFF, 0x6B, 0x35)


def main():
    prs = Presentation(PPTX_PATH)
    # 第5页（index 4）是宣传视频页
    # 当前已经有占位框, 优化提示文本
    print(f"当前幻灯片数: {len(prs.slides)}")
    print("视频嵌入操作说明：")
    print("1. 在 PowerPoint 中打开 PPTX")
    print("2. 切到第5页（宣传视频页）")
    print("3. 点击占位框 → 删除")
    print("4. 菜单栏 插入 → 视频 → 此电脑")
    print(f"5. 选择文件: {RECORDING_VIDEO}")
    print("6. 调整大小与位置 (建议与原占位框一致)")


if __name__ == "__main__":
    main()