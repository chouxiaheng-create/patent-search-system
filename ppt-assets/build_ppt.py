"""专利检索智能体 — 项目汇报PPT生成脚本.

共14页, 16:9 宽屏 (13.333 x 7.5 英寸), 包含:
- 封面、议程、背景、目标
- Remotion宣传视频占位 / 视频演示页
- 架构、核心能力、工作流、并行机制
- 看板截图、报告样例、质量感知
- 规划、收尾
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ============== 配色方案 ==============
COLOR_PRIMARY = RGBColor(0x38, 0x59, 0xFF)     # 主蓝
COLOR_DARK = RGBColor(0x1C, 0x20, 0x38)         # 标题深色
COLOR_ACCENT = RGBColor(0xFF, 0x6B, 0x35)       # 强调橙
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)     # 浅灰背景
COLOR_BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)      # 卡片白
COLOR_TEXT = RGBColor(0x2A, 0x2D, 0x3A)         # 正文
COLOR_TEXT_LIGHT = RGBColor(0x6B, 0x70, 0x80)   # 次要文字
COLOR_BORDER = RGBColor(0xE5, 0xE8, 0xF0)       # 边框灰
COLOR_GREEN = RGBColor(0x10, 0xB9, 0x81)        # 成功绿
COLOR_RED = RGBColor(0xEF, 0x44, 0x44)          # 错误红

# 资源路径
ASSETS_DIR = r"D:\Claude Code Files\Project_Patent search system_v1\ppt-assets"
RECORDING_VIDEO = r"D:\审协河南中心\【04】大模型应用小组-2026\专利检索智能体操作录屏.mp4"
KANBAN_IMAGE = os.path.join(ASSETS_DIR, "kanban_preview.png")
REPORT_P1 = os.path.join(ASSETS_DIR, "report_pages", "report_p1.png")
REPORT_P2 = os.path.join(ASSETS_DIR, "report_pages", "report_p2.png")
REPORT_P3 = os.path.join(ASSETS_DIR, "report_pages", "report_p3.png")
LOGOS_DIR = os.path.join(ASSETS_DIR, "logos")
LOGO_FILES = {
    "Kimi": os.path.join(LOGOS_DIR, "Kimi.png"),
    "智谱": os.path.join(LOGOS_DIR, "智谱.png"),
    "通义千问": os.path.join(LOGOS_DIR, "通义千问.png"),
    "DeepSeek": os.path.join(LOGOS_DIR, "DeepSeek.png"),
    "秘塔": os.path.join(LOGOS_DIR, "秘塔.png"),
    "OpenAI": os.path.join(LOGOS_DIR, "OpenAI.png"),
}

# 输出路径
OUTPUT_PPTX = r"D:\Claude Code Files\Project_Patent search system_v1\专利检索智能体汇报PPT.pptx"


# ============== 工具函数 ==============
def set_run_font(run, name="Microsoft YaHei", size=18, bold=False, color=COLOR_TEXT):
    """设置 run 字体 (中英文兼容)."""
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 设置东亚字体
    rPr = run._r.get_or_add_rPr()
    eastAsia = rPr.find(qn("a:ea"))
    if eastAsia is None:
        eastAsia = etree.SubElement(rPr, qn("a:ea"))
    eastAsia.set("typeface", "Microsoft YaHei")


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    """在幻灯片上添加文本框."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        set_run_font(run, size=font_size, bold=bold, color=color)
    return tb


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None,
             line_width=0.5, corner_radius=None):
    """添加矩形 (支持圆角)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.shadow.inherit = False
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape


def add_line(slide, x1, y1, x2, y2, color=COLOR_BORDER, width=1.0):
    """添加直线."""
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_bullets(slide, left, top, width, height, items, font_size=16,
                color=COLOR_TEXT, bullet_color=COLOR_PRIMARY, line_spacing=1.4):
    """添加带圆点项目列表."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        # 圆点
        run_dot = p.add_run()
        run_dot.text = "●  "
        set_run_font(run_dot, size=font_size, bold=True, color=bullet_color)
        # 内容
        run_text = p.add_run()
        run_text.text = item
        set_run_font(run_text, size=font_size, color=color)
    return tb


def add_page_header(slide, title, page_no, total=14):
    """添加页眉: 左侧标题, 右侧页码."""
    add_text_box(slide, 0.5, 0.3, 8, 0.5, title, font_size=20, bold=True, color=COLOR_DARK)
    add_text_box(slide, 12, 0.35, 0.8, 0.4, f"{page_no} / {total}",
                 font_size=12, color=COLOR_TEXT_LIGHT, align=PP_ALIGN.RIGHT)
    # 顶部分割线
    add_line(slide, 0.5, 0.85, 12.83, 0.85, color=COLOR_BORDER, width=1.0)


def add_page_footer(slide, text="专利检索智能体 · 项目阶段汇报"):
    """添加页脚."""
    add_text_box(slide, 0.5, 7.15, 10, 0.3, text, font_size=10, color=COLOR_TEXT_LIGHT)
    add_text_box(slide, 11, 7.15, 1.83, 0.3, "周夏恒", font_size=10,
                 color=COLOR_TEXT_LIGHT, align=PP_ALIGN.RIGHT)


def add_video_placeholder(slide, video_path, left, top, width, height, caption=""):
    """添加视频占位框 (由于python-pptx 1.0.2不支持add_movie, 我们使用占位+说明的方式).

    实际视频嵌入步骤: 打开PPTX后, 删除该占位框, 插入->视频->此电脑->选择mp4.
    """
    # 深色占位框
    add_rect(slide, left, top, width, height, fill_color=COLOR_DARK,
             line_color=COLOR_PRIMARY, line_width=2.0, corner_radius=0.02)
    # 标题
    add_text_box(slide, left, top + height/2 - 0.8, width, 0.5, "🎬 系统录屏视频",
                 font_size=28, bold=True, color=COLOR_BG_CARD, align=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + height/2 - 0.2, width, 0.4,
                 "专利检索智能体完整操作录屏（2分38秒）",
                 font_size=16, color=RGBColor(0xC0, 0xC8, 0xE0), align=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + height/2 + 0.3, width, 0.4,
                 f"文件：{os.path.basename(video_path)}",
                 font_size=12, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + height/2 + 0.7, width, 0.4,
                 "提示：PowerPoint 中可点击「插入 → 视频 → 此电脑」嵌入真实视频",
                 font_size=10, color=COLOR_TEXT_LIGHT, align=PP_ALIGN.CENTER)
    if caption:
        add_text_box(slide, left, top + height + 0.1, width, 0.4, caption,
                     font_size=14, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)


# ============== 各页内容 ==============
def page_01_cover(prs):
    """封面页."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白页
    # 背景: 深色块
    add_rect(slide, 0, 0, 13.333, 7.5, fill_color=COLOR_DARK)
    # 装饰蓝色条
    add_rect(slide, 0, 6.5, 13.333, 0.06, fill_color=COLOR_PRIMARY)
    # 主标题
    add_text_box(slide, 0.8, 2.3, 12, 1.2, "专利检索智能体",
                 font_size=54, bold=True, color=COLOR_BG_CARD, align=PP_ALIGN.LEFT)
    # 副标题
    add_text_box(slide, 0.8, 3.5, 12, 0.8,
                 "AI 辅助专利审查系统 — 方案设计与成果演示",
                 font_size=26, color=RGBColor(0xC8, 0xCE, 0xE0), align=PP_ALIGN.LEFT)
    # 副副标题
    add_text_box(slide, 0.8, 4.3, 12, 0.5,
                 "Multi-Model Parallel Search · Automated Prior-Art Retrieval · Traceable Reports",
                 font_size=14, color=RGBColor(0x88, 0x92, 0xB0), align=PP_ALIGN.LEFT)
    # 演讲人信息
    add_text_box(slide, 0.8, 5.6, 6, 0.4, "汇报人：周夏恒",
                 font_size=18, color=COLOR_BG_CARD, align=PP_ALIGN.LEFT)
    add_text_box(slide, 0.8, 6.0, 6, 0.4, "部门：医化部 · 大模型应用工作组",
                 font_size=16, color=RGBColor(0xA0, 0xA8, 0xC0), align=PP_ALIGN.LEFT)
    add_text_box(slide, 0.8, 6.4, 6, 0.4, "2026 年 6 月",
                 font_size=16, color=RGBColor(0xA0, 0xA8, 0xC0), align=PP_ALIGN.LEFT)
    # 右上角小标签
    add_rect(slide, 10.5, 0.6, 2.2, 0.5, fill_color=COLOR_PRIMARY, corner_radius=0.3)
    add_text_box(slide, 10.5, 0.65, 2.2, 0.4, "项目阶段汇报",
                 font_size=14, bold=True, color=COLOR_BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def page_02_agenda(prs):
    """目录页."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "汇报议程", 2)
    add_page_footer(slide)

    items = [
        ("01", "项目背景与目标", "传统检索痛点 · 三大目标"),
        ("02", "方案介绍", "系统架构 · 工作流程 · 核心能力"),
        ("03", "成果演示", "看板截图 · 报告样例 · 录屏展示"),
        ("04", "现状与规划", "完成度 · 后续方向"),
    ]
    y0 = 1.6
    h = 1.1
    gap = 0.2
    for i, (num, title, sub) in enumerate(items):
        y = y0 + i * (h + gap)
        # 编号圆块
        add_rect(slide, 0.8, y, 1.1, h, fill_color=COLOR_PRIMARY, corner_radius=0.15)
        add_text_box(slide, 0.8, y, 1.1, h, num, font_size=36, bold=True,
                     color=COLOR_BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 标题+副标题
        add_text_box(slide, 2.2, y + 0.15, 10, 0.55, title, font_size=24, bold=True, color=COLOR_DARK)
        add_text_box(slide, 2.2, y + 0.7, 10, 0.4, sub, font_size=15, color=COLOR_TEXT_LIGHT)


def page_03_background(prs):
    """项目背景: 传统检索痛点."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "项目背景：传统检索的痛点", 3)
    add_page_footer(slide)

    add_text_box(slide, 0.6, 1.1, 12, 0.6,
                 "在日常专利审查工作中，现有检索工具存在四类突出问题：",
                 font_size=18, color=COLOR_TEXT_LIGHT)

    pains = [
        ("01", "经验依赖", "现有工具高度依赖审查员个人经验，\n新手上手周期长、检索质量参差"),
        ("02", "覆盖有限", "单一引擎覆盖度有限，跨语种、\n跨领域文献容易遗漏"),
        ("03", "汇总繁重", "多份对比文献需人工逐条阅读、\n汇总去重，工作量极大"),
        ("04", "追溯困难", "检索结果缺乏结构化呈现，\n引用追溯不便"),
    ]
    y0 = 2.0
    h = 2.2
    w = 2.95
    gap = 0.15
    for i, (num, title, desc) in enumerate(pains):
        x = 0.5 + i * (w + gap)
        add_rect(slide, x, y0, w, h, fill_color=COLOR_BG_CARD, line_color=COLOR_BORDER, line_width=1.0, corner_radius=0.08)
        # 编号
        add_rect(slide, x + 0.2, y0 + 0.2, 0.7, 0.7, fill_color=COLOR_PRIMARY, corner_radius=0.5)
        add_text_box(slide, x + 0.2, y0 + 0.2, 0.7, 0.7, num, font_size=18, bold=True,
                     color=COLOR_BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 标题
        add_text_box(slide, x + 0.2, y0 + 1.0, w - 0.4, 0.5, title,
                     font_size=20, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
        # 描述
        add_text_box(slide, x + 0.2, y0 + 1.55, w - 0.4, h - 1.7, desc,
                     font_size=13, color=COLOR_TEXT, align=PP_ALIGN.CENTER, line_spacing=1.4)

    # 底部金句
    add_rect(slide, 0.5, 5.5, 12.33, 1.0, fill_color=COLOR_DARK, corner_radius=0.05)
    add_text_box(slide, 0.8, 5.55, 12, 0.5, "项目使命",
                 font_size=14, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.LEFT)
    add_text_box(slide, 0.8, 5.95, 12, 0.5,
                 "让审查员从机械检索中解放出来，专注于创造性判断",
                 font_size=20, bold=True, color=COLOR_BG_CARD, align=PP_ALIGN.LEFT)


def page_04_goals(prs):
    """项目目标."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "项目目标", 4)
    add_page_footer(slide)

    add_text_box(slide, 0.6, 1.1, 12, 0.6,
                 "围绕上述痛点，项目设定三大目标：",
                 font_size=18, color=COLOR_TEXT_LIGHT)

    goals = [
        ("提效", "EFFICIENCY", "单件专利的对比文件检索\n耗时大幅压缩", COLOR_PRIMARY),
        ("提质", "QUALITY", "多模型并行检索，覆盖度\n优于任何单一引擎", COLOR_GREEN),
        ("可追溯", "TRACEABILITY", "每条对比文献具备完整\n来源链路，便于审查引用", COLOR_ACCENT),
    ]
    y0 = 2.0
    h = 3.5
    w = 4.0
    gap = 0.25
    for i, (zh, en, desc, color) in enumerate(goals):
        x = 0.5 + i * (w + gap)
        add_rect(slide, x, y0, w, h, fill_color=COLOR_BG_CARD, line_color=color, line_width=2.0, corner_radius=0.05)
        # 顶部色条
        add_rect(slide, x, y0, w, 0.6, fill_color=color, corner_radius=0.05)
        add_text_box(slide, x, y0 + 0.1, w, 0.4, zh,
                     font_size=22, bold=True, color=COLOR_BG_CARD, align=PP_ALIGN.CENTER)
        # 英文
        add_text_box(slide, x, y0 + 0.8, w, 0.4, en,
                     font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        # 描述
        add_text_box(slide, x + 0.3, y0 + 1.6, w - 0.6, 1.5, desc,
                     font_size=16, color=COLOR_TEXT, align=PP_ALIGN.CENTER, line_spacing=1.5)

    add_text_box(slide, 0.5, 5.8, 12.33, 1.0,
                 "✓ 业务价值：减少审查员机械工作量    ✓ 技术价值：构建可扩展的多模型检索框架    ✓ 工程价值：全流程可追溯可审计",
                 font_size=15, color=COLOR_TEXT_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def page_05_promo_video(prs):
    """宣传视频/录屏展示页."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "系统宣传片：核心价值主张", 5)
    add_page_footer(slide)

    # 视频/录屏区
    video_left, video_top, video_w, video_h = 0.6, 1.1, 9.2, 5.7
    if os.path.exists(RECORDING_VIDEO):
        # 嵌入您的真实录屏 (使用占位框, 提示手动嵌入)
        add_video_placeholder(slide, RECORDING_VIDEO, video_left, video_top, video_w, video_h)
        video_caption = ""
    else:
        add_rect(slide, video_left, video_top, video_w, video_h,
                 fill_color=COLOR_BG_LIGHT, line_color=COLOR_BORDER, line_width=1.0)
        add_text_box(slide, video_left, video_top + video_h/2 - 0.3, video_w, 0.6,
                     "🎬 系统宣传视频", font_size=24, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
        video_caption = ""

    if video_caption:
        add_text_box(slide, video_left, video_top + video_h + 0.05, video_w, 0.4,
                     video_caption, font_size=13, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)

    # 右侧: 视频要点
    side_x = 10.1
    add_rect(slide, side_x, 1.1, 2.8, 5.7, fill_color=COLOR_DARK, corner_radius=0.05)
    add_text_box(slide, side_x + 0.2, 1.25, 2.5, 0.4, "视频要点",
                 font_size=16, bold=True, color=COLOR_BG_CARD)
    add_line(slide, side_x + 0.2, 1.75, side_x + 2.6, 1.75, color=COLOR_PRIMARY, width=1.5)
    key_points = [
        "登录与权限管理",
        "专利文件上传",
        "AI 自动解析",
        "多模型配置",
        "并行检索执行",
        "实时看板监控",
        "报告自动生成",
        "导出与评级",
    ]
    for i, p in enumerate(key_points):
        y = 1.95 + i * 0.55
        add_rect(slide, side_x + 0.25, y + 0.13, 0.15, 0.15, fill_color=COLOR_PRIMARY)
        add_text_box(slide, side_x + 0.5, y, 2.3, 0.45, p, font_size=14, color=COLOR_BG_CARD)


def page_06_architecture(prs):
    """系统架构."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "系统架构：三层 + 多模型适配层", 6)
    add_page_footer(slide)

    # 三层架构
    layers = [
        ("表现层", "Next.js 16 · React 19 · TypeScript", "用户界面 · 流程看板 · 报告查看 · 任务配置", COLOR_PRIMARY),
        ("业务层", "Node.js Worker · pg-boss 队列", "任务调度 · AI 调度 · 报告生成 · 失败重试", COLOR_GREEN),
        ("数据层", "Supabase (Postgres + Storage + Auth + Realtime)", "数据持久化 · 文件存储 · 身份认证 · 实时推送", COLOR_ACCENT),
    ]
    y = 1.1
    h = 1.3
    for i, (name, stack, func, color) in enumerate(layers):
        add_rect(slide, 0.5, y, 8.5, h, fill_color=COLOR_BG_CARD, line_color=color, line_width=2.0, corner_radius=0.05)
        # 左侧色块
        add_rect(slide, 0.5, y, 1.5, h, fill_color=color, corner_radius=0.05)
        add_text_box(slide, 0.5, y, 1.5, h, name, font_size=22, bold=True,
                     color=COLOR_BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, 2.2, y + 0.15, 6.7, 0.45, stack,
                     font_size=14, bold=True, color=color)
        add_text_box(slide, 2.2, y + 0.6, 6.7, 0.6, func,
                     font_size=13, color=COLOR_TEXT, line_spacing=1.4)
        y += h + 0.15

    # 右侧: 多模型适配层
    add_rect(slide, 9.3, 1.1, 3.5, 4.0, fill_color=COLOR_DARK, corner_radius=0.05)
    add_text_box(slide, 9.5, 1.25, 3.3, 0.4, "多 AI 引擎适配层",
                 font_size=16, bold=True, color=COLOR_BG_CARD)
    add_line(slide, 9.5, 1.75, 12.6, 1.75, color=COLOR_PRIMARY, width=1.5)

    models = ["Kimi K2.6", "智谱 GLM-5.1", "通义千问 Qwen", "DeepSeek", "秘塔 AI", "OpenAI 兼容"]
    for i, m in enumerate(models):
        yi = 1.95 + i * 0.5
        add_rect(slide, 9.5, yi, 3.1, 0.4, fill_color=RGBColor(0x2C, 0x33, 0x55), corner_radius=0.2)
        add_text_box(slide, 9.5, yi, 3.1, 0.4, m,
                     font_size=13, color=COLOR_BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 底部: 关键设计
    add_rect(slide, 0.5, 5.5, 12.33, 1.3, fill_color=COLOR_BG_LIGHT, corner_radius=0.05)
    add_text_box(slide, 0.7, 5.6, 12, 0.4, "关键设计原则",
                 font_size=14, bold=True, color=COLOR_PRIMARY)
    add_text_box(slide, 0.7, 5.95, 12, 0.85,
                 "● 前后端分离：Worker 独立进程，通过 pg-boss 队列解耦    ● 实时同步：基于 Supabase Realtime 推送，前端无轮询\n"
                 "● 适配器工厂：各厂商差异封装在 AIAdapter 之后，新增模型零成本接入    ● 失败安全：指数退避重试 + 20 分钟全局超时 + 支持取消",
                 font_size=12, color=COLOR_TEXT, line_spacing=1.5)


def page_07_capabilities(prs):
    """六大核心能力 (含模型LOGO)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "六大核心能力", 7)
    add_page_footer(slide)

    # 顶部: 6个LOGO横排
    add_text_box(slide, 0.5, 1.05, 12, 0.4, "已集成 6 大 AI 引擎（可插拔适配）",
                 font_size=14, bold=True, color=COLOR_PRIMARY)
    logo_w = 1.85
    logo_h = 0.9
    logo_y = 1.5
    total_w = 6 * logo_w + 5 * 0.1
    start_x = (13.333 - total_w) / 2
    for i, (name, path) in enumerate(LOGO_FILES.items()):
        x = start_x + i * (logo_w + 0.1)
        if os.path.exists(path):
            slide.shapes.add_picture(path, Inches(x), Inches(logo_y),
                                     width=Inches(logo_w), height=Inches(logo_h))
        else:
            add_rect(slide, x, logo_y, logo_w, logo_h, fill_color=COLOR_BG_LIGHT, line_color=COLOR_BORDER, line_width=1.0)
            add_text_box(slide, x, logo_y, logo_w, logo_h, name,
                         font_size=16, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 下方: 6个能力卡片 (3x2)
    add_text_box(slide, 0.5, 2.65, 12, 0.4, "六大核心能力",
                 font_size=14, bold=True, color=COLOR_PRIMARY)
    caps = [
        ("多模型并行检索", "N 个模型 × M 个策略笛卡尔积\n并行执行，覆盖度大幅提升"),
        ("可插拔 AI 适配器", "统一接口屏蔽各厂商差异\n新增模型零成本接入"),
        ("可视化流程看板", "React Flow 实时呈现\n解析→检索→生成全链路"),
        ("智能去重与排序", "URL + 标题归一化去重\n独立报告模型 Top-N 选优"),
        ("质量感知与告警", "缺失字段自动标注\n低质量结果直接过滤"),
        ("可追溯审计", "来源平台/策略/任务 ID\n支持 GB/T 7714 引用"),
    ]
    y0 = 3.1
    h = 1.85
    w = 4.0
    gap_x = 0.15
    gap_y = 0.15
    for i, (name, desc) in enumerate(caps):
        col = i % 3
        row = i // 3
        x = 0.5 + col * (w + gap_x)
        y = y0 + row * (h + gap_y)
        add_rect(slide, x, y, w, h, fill_color=COLOR_BG_CARD, line_color=COLOR_BORDER, line_width=1.0, corner_radius=0.05)
        # 编号
        add_rect(slide, x + 0.15, y + 0.15, 0.5, 0.5, fill_color=COLOR_PRIMARY, corner_radius=0.4)
        add_text_box(slide, x + 0.15, y + 0.15, 0.5, 0.5, str(i + 1),
                     font_size=16, bold=True, color=COLOR_BG_CARD,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 标题
        add_text_box(slide, x + 0.75, y + 0.2, w - 0.9, 0.45, name,
                     font_size=15, bold=True, color=COLOR_DARK)
        # 描述
        add_text_box(slide, x + 0.2, y + 0.85, w - 0.4, h - 1.0, desc,
                     font_size=12, color=COLOR_TEXT, line_spacing=1.4)


def page_08_workflow(prs):
    """工作流程."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "审查员工作流：六步完成检索", 8)
    add_page_footer(slide)

    # 流程: 横向4步, 加上看板与报告
    add_text_box(slide, 0.5, 1.05, 12, 0.4, "核心工作流（从上传到报告导出）",
                 font_size=14, bold=True, color=COLOR_PRIMARY)

    steps = [
        ("1", "上传专利", "PDF / DOCX /\nXLSX / TXT"),
        ("2", "AI 解析", "技术主题 · 申请人\n权利要求 · 核心发明点"),
        ("3", "配置检索", "模型 × 策略矩阵\n可微调默认值"),
        ("4", "启动并行", "后台自动执行\nN×M 子任务并行"),
        ("5", "流程看板", "React Flow 实时\n状态可视化"),
        ("6", "生成报告", "Top-N 自动选优\n支持 MD/DOCX 导出"),
    ]
    y0 = 1.7
    h = 1.4
    w = 1.95
    gap = 0.1
    for i, (num, name, desc) in enumerate(steps):
        x = 0.5 + i * (w + gap)
        add_rect(slide, x, y0, w, h, fill_color=COLOR_BG_CARD, line_color=COLOR_PRIMARY, line_width=1.5, corner_radius=0.08)
        # 编号圆
        add_rect(slide, x + 0.1, y0 + 0.1, 0.55, 0.55, fill_color=COLOR_PRIMARY, corner_radius=0.5)
        add_text_box(slide, x + 0.1, y0 + 0.1, 0.55, 0.55, num,
                     font_size=18, bold=True, color=COLOR_BG_CARD,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 名称
        add_text_box(slide, x, y0 + 0.75, w, 0.35, name,
                     font_size=14, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
        # 描述
        add_text_box(slide, x + 0.05, y0 + 1.1, w - 0.1, 0.3, desc,
                     font_size=10, color=COLOR_TEXT_LIGHT, align=PP_ALIGN.CENTER, line_spacing=1.3)
        # 箭头
        if i < len(steps) - 1:
            arr_x = x + w + 0.005
            arr_y = y0 + h / 2
            add_text_box(slide, arr_x, arr_y - 0.15, 0.1, 0.3, "→",
                         font_size=18, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)

    # 下方: 三步关键能力
    add_text_box(slide, 0.5, 3.5, 12, 0.4, "每步背后的关键能力",
                 font_size=14, bold=True, color=COLOR_PRIMARY)
    detail = [
        ("灵活的文件处理", "PDF/DOCX/XLSX/TXT 全格式支持，自动 OCR 与结构化抽取"),
        ("智能 Prompt 编辑", "解析与检索的 Prompt 可视化编辑，支持按用户经验定制"),
        ("人机协作闭环", "报告支持人工评级、备注、修改，形成知识沉淀"),
    ]
    y0 = 3.95
    h = 0.95
    w = 4.0
    gap = 0.15
    for i, (name, desc) in enumerate(detail):
        x = 0.5 + i * (w + gap)
        add_rect(slide, x, y0, w, h, fill_color=COLOR_BG_LIGHT, corner_radius=0.05)
        add_text_box(slide, x + 0.2, y0 + 0.1, w - 0.4, 0.4, name,
                     font_size=14, bold=True, color=COLOR_PRIMARY)
        add_text_box(slide, x + 0.2, y0 + 0.5, w - 0.4, h - 0.55, desc,
                     font_size=12, color=COLOR_TEXT, line_spacing=1.4)

    # 登录页截图锚点
    add_text_box(slide, 0.5, 5.2, 12, 0.4, "登录界面（系统入口）",
                 font_size=14, bold=True, color=COLOR_PRIMARY)
    login_path = os.path.join(ASSETS_DIR, "login-preview.png")
    if os.path.exists(login_path):
        # 等比缩放
        from PIL import Image as PImage
        img = PImage.open(login_path)
        ratio = img.height / img.width
        w_pic = 4.0
        h_pic = w_pic * ratio
        if h_pic > 1.7:
            h_pic = 1.7
            w_pic = h_pic / ratio
        slide.shapes.add_picture(login_path, Inches((13.333 - w_pic) / 2), Inches(5.65),
                                 width=Inches(w_pic), height=Inches(h_pic))


def page_09_parallel(prs):
    """核心机制: 笛卡尔积并行."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "核心机制：模型 × 策略 笛卡尔积并行", 9)
    add_page_footer(slide)

    add_text_box(slide, 0.5, 1.1, 12, 0.5,
                 "N 个模型 × M 个策略 → 自动展开 N×M 个并行子任务",
                 font_size=16, bold=True, color=COLOR_DARK)

    # 矩阵示意
    add_text_box(slide, 0.5, 1.85, 12, 0.4, "并行执行矩阵",
                 font_size=14, bold=True, color=COLOR_PRIMARY)

    models = ["智谱GLM-5.1", "秘塔AI", "Kimi K2.6"]
    strategies = ["追踪检索", "主要技术方案步骤检索", "发明构思检索"]
    matrix_x = 0.5
    matrix_y = 2.4
    cell_w = 2.2
    cell_h = 0.6
    # 顶部标题
    for j, s in enumerate(strategies):
        x = matrix_x + (j + 1) * cell_w
        add_rect(slide, x, matrix_y, cell_w, cell_h, fill_color=COLOR_PRIMARY)
        add_text_box(slide, x, matrix_y, cell_w, cell_h, s,
                     font_size=12, bold=True, color=COLOR_BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 行标题 + 单元格
    for i, m in enumerate(models):
        y = matrix_y + (i + 1) * cell_h
        add_rect(slide, matrix_x, y, cell_w, cell_h, fill_color=COLOR_DARK)
        add_text_box(slide, matrix_x, y, cell_w, cell_h, m,
                     font_size=12, bold=True, color=COLOR_BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        for j in range(len(strategies)):
            x = matrix_x + (j + 1) * cell_w
            add_rect(slide, x, y, cell_w, cell_h, fill_color=COLOR_BG_CARD, line_color=COLOR_BORDER, line_width=0.5)
            task_id = f"任务{i+1}{chr(65+j)}"
            add_text_box(slide, x, y, cell_w, cell_h, task_id,
                         font_size=12, bold=True, color=COLOR_GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 右侧说明
    add_rect(slide, 8.5, 1.85, 4.3, 3.0, fill_color=COLOR_BG_LIGHT, corner_radius=0.05)
    add_text_box(slide, 8.7, 1.95, 4, 0.4, "关键设计",
                 font_size=14, bold=True, color=COLOR_PRIMARY)
    design_points = [
        "● 单模型并发上限 2，避免触发厂商限流",
        "● 失败任务自动重试（指数退避策略）",
        "● 全局 20 分钟超时保护",
        "● 全程支持审查员取消操作",
        "● 子任务完成后自动汇入统一去重与排序模块",
        "● 由独立的\"报告模型\"完成 Top-N 选优",
    ]
    add_bullets(slide, 8.7, 2.4, 4.0, 2.4, design_points, font_size=12, line_spacing=1.4)

    # 底部: 价值
    add_rect(slide, 0.5, 5.5, 12.33, 1.3, fill_color=COLOR_DARK, corner_radius=0.05)
    add_text_box(slide, 0.8, 5.6, 12, 0.4, "为什么是\"集思广益\"？",
                 font_size=14, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, 0.8, 6.0, 12, 0.85,
                 "不同 AI 引擎的训练数据与检索偏好不同，单一引擎存在系统性盲区；多模型并行执行后再统一去重与排序，\n"
                 "既扩大了覆盖广度（更多候选文献），又通过共识机制提升结果质量（多源印证）。",
                 font_size=13, color=COLOR_BG_CARD, line_spacing=1.5)


def page_10_kanban(prs):
    """流程看板截图."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "成果演示一：实时流程看板", 10)
    add_page_footer(slide)

    add_text_box(slide, 0.5, 1.1, 12, 0.5,
                 "看板基于 React Flow 实现，每个子任务对应一个可视化节点，状态实时反映后端执行进度",
                 font_size=14, color=COLOR_TEXT_LIGHT)

    if os.path.exists(KANBAN_IMAGE):
        # 等比缩放至最大可用空间
        from PIL import Image as PImage
        img = PImage.open(KANBAN_IMAGE)
        ratio = img.width / img.height
        max_w = 12.0
        max_h = 4.5
        if ratio > max_w / max_h:
            w_pic = max_w
            h_pic = w_pic / ratio
        else:
            h_pic = max_h
            w_pic = h_pic * ratio
        pic_x = (13.333 - w_pic) / 2
        slide.shapes.add_picture(KANBAN_IMAGE, Inches(pic_x), Inches(1.7),
                                 width=Inches(w_pic), height=Inches(h_pic))
        # 提示替换
        add_text_box(slide, 0.5, 1.7 + h_pic + 0.1, 12, 0.4,
                     "💡 提示：此为自动生成的示意图，建议替换为您的真实看板截图（消息中已上传）",
                     font_size=11, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    else:
        add_rect(slide, 0.5, 1.7, 12.33, 4.5, fill_color=COLOR_BG_LIGHT, line_color=COLOR_BORDER, line_width=1.0)
        add_text_box(slide, 0.5, 3.5, 12.33, 0.6, "看板截图占位",
                     font_size=18, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)

    # 底部: 看板亮点
    add_rect(slide, 0.5, 6.5, 12.33, 0.6, fill_color=COLOR_BG_LIGHT, corner_radius=0.05)
    add_text_box(slide, 0.7, 6.55, 12, 0.5,
                 "亮点：四态可视化（待执行/执行中/成功/失败） · 支持取消 · 排队位置可见 · Supabase Realtime 推送，零轮询",
                 font_size=12, color=COLOR_PRIMARY, anchor=MSO_ANCHOR.MIDDLE)


def page_11_report_sample(prs):
    """报告样例展示."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "成果演示二：报告样例（前 3 篇）", 11)
    add_page_footer(slide)

    add_text_box(slide, 0.5, 1.0, 12, 0.4,
                 "真实运行样例：待审专利 CN118429689A · 对比文献 10 篇（Google Patents + arXiv）",
                 font_size=13, color=COLOR_TEXT_LIGHT)

    if os.path.exists(REPORT_P1):
        from PIL import Image as PImage
        img = PImage.open(REPORT_P1)
        ratio = img.width / img.height
        max_w = 8.5
        max_h = 5.5
        if ratio > max_w / max_h:
            w_pic = max_w
            h_pic = w_pic / ratio
        else:
            h_pic = max_h
            w_pic = h_pic * ratio
        pic_x = 0.5
        slide.shapes.add_picture(REPORT_P1, Inches(pic_x), Inches(1.45),
                                 width=Inches(w_pic), height=Inches(h_pic))
    else:
        add_rect(slide, 0.5, 1.45, 8.5, 5.5, fill_color=COLOR_BG_LIGHT, line_color=COLOR_BORDER)
        add_text_box(slide, 0.5, 3.7, 8.5, 0.6, "报告样例占位",
                     font_size=18, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)

    # 右侧: 报告特性
    side_x = 9.3
    add_rect(slide, side_x, 1.45, 3.5, 5.5, fill_color=COLOR_DARK, corner_radius=0.05)
    add_text_box(slide, side_x + 0.2, 1.6, 3.3, 0.4, "报告关键特性",
                 font_size=15, bold=True, color=COLOR_BG_CARD)
    add_line(slide, side_x + 0.2, 2.05, side_x + 3.4, 2.05, color=COLOR_PRIMARY, width=1.5)
    feats = [
        ("结构化字段", "标题/来源/作者/\n日期/链接/描述"),
        ("来源链路", "标注来源平台与\n检索策略（来源×策略）"),
        ("可追溯", "每篇文献携带\n任务 ID，便于复核"),
        ("标准化引用", "支持 GB/T 7714\n引用格式导出"),
        ("多格式导出", "Markdown + DOCX\n两种格式可选"),
    ]
    y0 = 2.2
    h = 0.95
    for i, (name, desc) in enumerate(feats):
        y = y0 + i * h
        add_rect(slide, side_x + 0.2, y, 0.1, h - 0.15, fill_color=COLOR_PRIMARY)
        add_text_box(slide, side_x + 0.4, y, 3.0, 0.3, name,
                     font_size=13, bold=True, color=COLOR_BG_CARD)
        add_text_box(slide, side_x + 0.4, y + 0.3, 3.0, 0.55, desc,
                     font_size=11, color=RGBColor(0xC0, 0xC8, 0xE0), line_spacing=1.3)


def page_12_quality(prs):
    """质量感知能力."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "成果演示三：质量感知与可追溯性", 12)
    add_page_footer(slide)

    add_text_box(slide, 0.5, 1.0, 12, 0.4,
                 "系统对每篇对比文献进行质量评估，自动暴露数据缺陷",
                 font_size=14, color=COLOR_TEXT_LIGHT)

    # 左侧: 质量规则
    add_rect(slide, 0.5, 1.5, 6.2, 5.0, fill_color=COLOR_BG_CARD, line_color=COLOR_BORDER, line_width=1.0, corner_radius=0.05)
    add_text_box(slide, 0.7, 1.6, 6.0, 0.4, "质量规则",
                 font_size=15, bold=True, color=COLOR_PRIMARY)
    add_line(slide, 0.7, 2.05, 6.5, 2.05, color=COLOR_BORDER, width=1.0)
    rules = [
        ("缺失字段告警", "若日期、链接、描述缺失\n系统如实标注，不做虚假填充"),
        ("作者未知标注", "若作者字段为\"未知\"\n保留原始状态而非编造"),
        ("质量分过滤", "每篇文献计算质量分 (0-100)\n低于阈值 (默认 50) 直接过滤"),
        ("疑似重复识别", "同族专利识别为相关\n但保留为独立条目供审查员判断"),
    ]
    y0 = 2.2
    h = 1.05
    for i, (name, desc) in enumerate(rules):
        y = y0 + i * h
        add_rect(slide, 0.7, y, 0.1, h - 0.15, fill_color=COLOR_ACCENT)
        add_text_box(slide, 0.9, y, 5.5, 0.4, name,
                     font_size=13, bold=True, color=COLOR_DARK)
        add_text_box(slide, 0.9, y + 0.4, 5.5, 0.6, desc,
                     font_size=11, color=COLOR_TEXT, line_spacing=1.4)

    # 右侧: 来源链路示例
    add_rect(slide, 7.0, 1.5, 5.8, 5.0, fill_color=COLOR_BG_CARD, line_color=COLOR_BORDER, line_width=1.0, corner_radius=0.05)
    add_text_box(slide, 7.2, 1.6, 5.5, 0.4, "来源链路（样例）",
                 font_size=15, bold=True, color=COLOR_PRIMARY)
    add_line(slide, 7.2, 2.05, 13.0, 2.05, color=COLOR_BORDER, width=1.0)
    add_text_box(slide, 7.2, 2.2, 5.5, 0.4, "每篇对比文献的来源标注格式：",
                 font_size=12, color=COLOR_TEXT_LIGHT)

    # 样例条目
    samples = [
        ("DeepSeek × 追踪检索", "1. 基于注意力频域生成对抗网络的对抗样本生成方法与装置", "CN114510724A"),
        ("智谱GLM-5.1 × 发明构思检索", "2. Low-frequency adversarial perturbations", "arXiv:2009.02596"),
        ("智谱GLM-5.1 × 发明构思检索", "3. AdvGAN: Generating Adversarial Examples", "arXiv:1801.02610"),
        ("智谱GLM-5.1 × 追踪检索", "5. 基于注意力频域GAN的对抗性仿真攻击方法与装置", "CN115906455A"),
        ("秘塔AI × 发明构思检索", "6. 文献1（标题缺失）", "—"),
    ]
    y0 = 2.7
    h = 0.7
    for i, (src, title, no) in enumerate(samples):
        y = y0 + i * h
        add_rect(slide, 7.2, y, 5.4, h - 0.05, fill_color=COLOR_BG_LIGHT, corner_radius=0.03)
        add_text_box(slide, 7.3, y + 0.05, 5.0, 0.3, f"来源：{src}",
                     font_size=10, color=COLOR_PRIMARY, bold=True)
        add_text_box(slide, 7.3, y + 0.32, 3.5, 0.3, title,
                     font_size=10, color=COLOR_TEXT)
        add_text_box(slide, 11.0, y + 0.32, 1.6, 0.3, no,
                     font_size=9, color=COLOR_TEXT_LIGHT, align=PP_ALIGN.RIGHT)


def page_13_roadmap(prs):
    """现状与规划."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(slide, "项目现状与后续规划", 13)
    add_page_footer(slide)

    # 现状
    add_rect(slide, 0.5, 1.0, 6.0, 5.5, fill_color=COLOR_BG_CARD, line_color=COLOR_BORDER, line_width=1.0, corner_radius=0.05)
    add_text_box(slide, 0.7, 1.1, 5.6, 0.4, "✓ 已完成能力",
                 font_size=16, bold=True, color=COLOR_GREEN)
    add_line(slide, 0.7, 1.55, 6.3, 1.55, color=COLOR_GREEN, width=2.0)
    done_items = [
        "检索流程全链路闭环（上传→解析→检索→报告→导出）",
        "6 大 AI 模型适配与并行调度",
        "React Flow 实时流程看板",
        "智能去重、Top-N 选优、质量感知",
        "GB/T 7714 引用格式输出",
        "Mock 模式（无 API Key 即可演示）",
    ]
    add_bullets(slide, 0.7, 1.7, 5.6, 4.0, done_items, font_size=13, bullet_color=COLOR_GREEN, line_spacing=1.6)

    # 规划
    add_rect(slide, 6.83, 1.0, 6.0, 5.5, fill_color=COLOR_BG_CARD, line_color=COLOR_BORDER, line_width=1.0, corner_radius=0.05)
    add_text_box(slide, 7.03, 1.1, 5.6, 0.4, "→ 后续规划（待领导指示）",
                 font_size=16, bold=True, color=COLOR_PRIMARY)
    add_line(slide, 7.03, 1.55, 12.63, 1.55, color=COLOR_PRIMARY, width=2.0)

    # 方向A
    add_rect(slide, 7.03, 1.7, 5.6, 2.2, fill_color=COLOR_BG_LIGHT, corner_radius=0.05)
    add_text_box(slide, 7.2, 1.8, 5.4, 0.4, "方向 A：检索策略深度优化",
                 font_size=13, bold=True, color=COLOR_DARK)
    a_items = [
        "● IPC 分类号联动",
        "● 引文网络图谱",
        "● 同族专利智能识别",
    ]
    add_bullets(slide, 7.2, 2.25, 5.4, 1.6, a_items, font_size=11, line_spacing=1.5)

    # 方向B
    add_rect(slide, 7.03, 4.0, 5.6, 2.2, fill_color=COLOR_BG_LIGHT, corner_radius=0.05)
    add_text_box(slide, 7.2, 4.1, 5.4, 0.4, "方向 B：报告系统增强",
                 font_size=13, bold=True, color=COLOR_DARK)
    b_items = [
        "● 审查建议自动生成",
        "● 权利要求对比表自动填充",
        "● 审查员知识库沉淀",
    ]
    add_bullets(slide, 7.2, 4.55, 5.4, 1.6, b_items, font_size=11, line_spacing=1.5)


def page_14_ending(prs):
    """结束页."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, fill_color=COLOR_DARK)
    add_rect(slide, 0, 3.2, 13.333, 0.04, fill_color=COLOR_PRIMARY)
    add_text_box(slide, 0.5, 1.5, 12.33, 1.0, "谢谢聆听",
                 font_size=60, bold=True, color=COLOR_BG_CARD, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 2.6, 12.33, 0.5, "Thank You for Your Attention",
                 font_size=20, color=RGBColor(0xC0, 0xC8, 0xE0), align=PP_ALIGN.CENTER)

    # 三大目标回顾
    add_text_box(slide, 0.5, 3.6, 12.33, 0.5, "三大目标回顾",
                 font_size=18, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    goals = [("提效", COLOR_PRIMARY), ("提质", COLOR_GREEN), ("可追溯", COLOR_ACCENT)]
    total_w = 3 * 2.5 + 2 * 0.5
    start_x = (13.333 - total_w) / 2
    for i, (g, c) in enumerate(goals):
        x = start_x + i * (2.5 + 0.5)
        add_rect(slide, x, 4.2, 2.5, 0.7, fill_color=c, corner_radius=0.1)
        add_text_box(slide, x, 4.2, 2.5, 0.7, g, font_size=22, bold=True,
                     color=COLOR_BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 价值主张
    add_text_box(slide, 0.5, 5.4, 12.33, 0.5, "让审查员从机械检索中解放出来，专注于创造性判断",
                 font_size=22, color=COLOR_BG_CARD, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0.5, 6.1, 12.33, 0.5, "—— 专利检索智能体项目愿景",
                 font_size=14, color=RGBColor(0x88, 0x92, 0xB0), align=PP_ALIGN.CENTER)

    # Q&A
    add_text_box(slide, 0.5, 6.6, 12.33, 0.5, "欢迎各位领导指正与提问",
                 font_size=16, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)


# ============== 主函数 ==============
def main():
    print("=" * 60)
    print("开始生成PPT: 专利检索智能体项目汇报")
    print("=" * 60)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        ("封面", page_01_cover),
        ("议程", page_02_agenda),
        ("背景", page_03_background),
        ("目标", page_04_goals),
        ("宣传视频", page_05_promo_video),
        ("架构", page_06_architecture),
        ("核心能力", page_07_capabilities),
        ("工作流", page_08_workflow),
        ("并行机制", page_09_parallel),
        ("看板截图", page_10_kanban),
        ("报告样例", page_11_report_sample),
        ("质量感知", page_12_quality),
        ("现状规划", page_13_roadmap),
        ("结束页", page_14_ending),
    ]
    for i, (name, fn) in enumerate(builders, 1):
        print(f"  [{i:02d}/14] {name} ...")
        fn(prs)
    print()
    prs.save(OUTPUT_PPTX)
    print(f"✓ PPT 已生成: {OUTPUT_PPTX}")
    print(f"  文件大小: {os.path.getsize(OUTPUT_PPTX) / 1024:.1f} KB")


if __name__ == "__main__":
    main()